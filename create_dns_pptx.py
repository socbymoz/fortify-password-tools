import os, io
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import numpy as np

plt.rcParams.update({
    'font.family': 'sans-serif',
    'font.sans-serif': ['Segoe UI', 'DejaVu Sans', 'Arial'],
    'text.color': '#FFFFFF',
    'axes.facecolor': '#1E1E2E',
    'figure.facecolor': '#1E1E2E',
    'axes.edgecolor': '#333',
    'axes.labelcolor': '#FFFFFF',
    'xtick.color': '#CCCCCC',
    'ytick.color': '#CCCCCC',
})

W = 12.0
H = 6.6
ACCENT = '#00A8E8'
WHITE  = '#FFFFFF'
LGRAY  = '#CCCCCC'
YELLOW = '#FFD700'
GREEN  = '#4CAF50'
ORANGE = '#FF9800'
RED    = '#F44336'
CARD   = '#2A2A3E'

def new_fig():
    fig, ax = plt.subplots(figsize=(W, H))
    ax.set_xlim(0, W); ax.set_ylim(0, H)
    ax.axis('off')
    return fig, ax

def card(ax, x, y, w, h, color=ACCENT, lw=2):
    bx = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.08",
                        facecolor=CARD, edgecolor=color, linewidth=lw)
    ax.add_patch(bx)
    return bx

def txt(ax, x, y, text, size=13, color=WHITE, bold=False, ha='left', va='center', alpha=1):
    ax.text(x, y, text, fontsize=size, color=color, fontweight='bold' if bold else 'normal',
            ha=ha, va=va, alpha=alpha)

def title_bar(ax, text):
    card(ax, 0.2, H-0.75, W-0.4, 0.55, ACCENT, 2)
    txt(ax, 0.5, H-0.47, text, 22, WHITE, True, 'left', 'center')

def fig_to_bytes(fig, dpi=130):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight', facecolor=fig.get_facecolor(), edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf

# ═══ SLIDE 1 — Title ═══
def s1_title():
    fig, ax = new_fig()
    txt(ax, W/2, H-1.0, "DNS Resolution", 44, WHITE, True, 'center', 'center')
    txt(ax, W/2, H-1.9, "How the Internet Finds Its Way", 26, ACCENT, False, 'center', 'center')
    txt(ax, W/2, H-2.5, "Deep Dive into the Domain Name System", 16, LGRAY, False, 'center', 'center')
    # Phonebook visual
    card(ax, 0.8, 1.2, 4.5, 2.2, ACCENT, 3)
    txt(ax, 3.05, 2.3, "www.example.com", 18, WHITE, True, 'center', 'center')
    txt(ax, 3.05, 1.7, "Domain Name", 12, LGRAY, False, 'center', 'center')
    # Arrow
    ax.annotate('', xy=(7.8, 2.3), xytext=(5.6, 2.3),
                arrowprops=dict(arrowstyle='->', color=ACCENT, lw=4))
    txt(ax, 6.7, 2.9, "DNS Resolution", 13, YELLOW, True, 'center', 'center')
    # IP box
    card(ax, 8.1, 1.2, 3.5, 2.2, GREEN, 3)
    txt(ax, 9.85, 2.3, "192.0.2.1", 18, GREEN, True, 'center', 'center')
    txt(ax, 9.85, 1.7, "IP Address", 12, LGRAY, False, 'center', 'center')
    # Bottom info
    txt(ax, W/2, 0.4, "www.student-login.local", 11, LGRAY, False, 'center', 'center')
    return fig

# ═══ SLIDE 2 — What is DNS ═══
def s2_what():
    fig, ax = new_fig()
    title_bar(ax, "What is DNS?")
    txt(ax, 0.5, H-1.5, "The Domain Name System (DNS) is the phonebook of the Internet.", 18, WHITE, True)
    items = [
        "Translates human-friendly domain names into machine-readable IP addresses",
        "Eliminates the need to memorize numeric IP addresses",
        "A distributed, hierarchical database with no single point of failure",
        "One of the oldest and most critical components of modern networking",
        "Operates on UDP port 53 (with TCP fallback for large responses)",
    ]
    for i, item in enumerate(items):
        y = H - 2.3 - i * 0.55
        ax.plot(0.6, y, 'o', color=ACCENT, markersize=6)
        txt(ax, 0.9, y, item, 14, LGRAY)

    # Visual - DNS translation
    card(ax, 0.6, 0.5, 5.0, 1.5, ACCENT, 2)
    txt(ax, 1.0, 1.5, "Type:  google.com", 15, WHITE, True)
    txt(ax, 1.0, 1.0, "Get:   142.250.190.46", 15, GREEN, True)

    card(ax, 6.0, 0.5, 5.6, 1.5, YELLOW, 2)
    txt(ax, 6.4, 1.5, "Type:  github.com", 15, WHITE, True)
    txt(ax, 6.4, 1.0, "Get:   140.82.121.3", 15, GREEN, True)
    return fig

# ═══ SLIDE 3 — Why DNS Matters ═══
def s3_why():
    fig, ax = new_fig()
    title_bar(ax, "Why DNS Matters")
    items = [
        ("User Experience", "Names are easier to remember than numbers", ACCENT),
        ("Performance", "Caching at every level reduces lookup latency", GREEN),
        ("Security", "DNSSEC, DNS-over-HTTPS, DNS-over-TLS", ORANGE),
        ("Load Balancing", "Round-robin DNS distributes traffic across servers", YELLOW),
        ("Filtering", "Block malicious domains at the resolver level", RED),
        ("Service Discovery", "SRV records locate services automatically", GREEN),
        ("Resilience", "Anycast routing and redundant nameservers", ACCENT),
    ]
    for i, (title, desc, color) in enumerate(items):
        y = H - 1.7 - i * 0.6
        card(ax, 0.4, y-0.2, 11.2, 0.45, color, 1.5)
        txt(ax, 0.7, y, title, 15, color, True)
        txt(ax, 3.5, y, desc, 13, LGRAY)
    return fig

# ═══ SLIDE 4 — Full Pipeline Overview ═══
def s4_pipeline():
    fig, ax = new_fig()
    title_bar(ax, "DNS Resolution — Full Pipeline")
    steps = [
        (1, "Browser\nCache", ACCENT),
        (2, "OS / Hosts\nFile", ACCENT),
        (3, "Recursive\nResolver", ACCENT),
        (4, "Root\nServer", ORANGE),
        (5, "TLD\nServer", ORANGE),
        (6, "Authoritative\nServer", GREEN),
    ]
    xs = [0.3, 2.2, 4.1, 6.0, 7.9, 9.8]
    for (num, label, color), x in zip(steps, xs):
        card(ax, x, 2.8, 1.7, 1.8, color, 3)
        txt(ax, x+0.85, 4.05, str(num), 13, color, True, 'center')
        txt(ax, x+0.85, 3.6, label, 15, WHITE, True, 'center', 'center')
        if num < 6:
            ax.annotate('', xy=(x+1.9, 3.7), xytext=(x+1.7, 3.7),
                        arrowprops=dict(arrowstyle='->', color=ACCENT, lw=3))

    # Client
    card(ax, 3.5, 5.3, 2.5, 0.6, YELLOW, 2)
    txt(ax, 4.75, 5.6, "User Types URL", 14, YELLOW, True, 'center', 'center')
    ax.annotate('', xy=(3.75, 4.6), xytext=(4.75, 5.3),
                arrowprops=dict(arrowstyle='->', color=LGRAY, lw=1.5))

    txt(ax, 6.0, 2.0, "Each step delegates to the next until the IP address is found", 11, LGRAY, False, 'center', 'center')
    return fig

# ═══ SLIDE 5 — Browser Cache ═══
def s5_browser():
    fig, ax = new_fig()
    title_bar(ax, "Step 1: Browser DNS Cache")
    card(ax, 0.4, 4.5, 5.5, 1.6, ACCENT, 2.5)
    txt(ax, 3.15, 5.8, "Browser DNS Cache (In-Memory)", 16, ACCENT, True, 'center')
    txt(ax, 3.15, 5.3, "Stores recent DNS lookups per browser session", 12, LGRAY, False, 'center')
    txt(ax, 3.15, 4.9, "TTL values determine how long entries persist", 12, LGRAY, False, 'center')

    # Example cache entries
    card(ax, 6.2, 4.5, 5.4, 1.6, GREEN, 1.5)
    txt(ax, 6.5, 5.8, "Example Cache Entries", 14, GREEN, True)
    txt(ax, 6.5, 5.3, "google.com  →  142.250.190.46", 12, WHITE)
    txt(ax, 6.5, 4.9, "github.com  →  140.82.121.3", 12, WHITE)

    # Decision flow
    card(ax, 0.4, 1.8, 5.5, 1.8, GREEN, 2)
    txt(ax, 3.15, 3.0, "Found in cache?", 16, GREEN, True, 'center')
    txt(ax, 3.15, 2.4, "Immediate resolution (sub-millisecond)", 13, WHITE, False, 'center')

    card(ax, 6.2, 1.8, 5.4, 1.8, RED, 2)
    txt(ax, 8.9, 3.0, "Not found?", 16, RED, True, 'center')
    txt(ax, 8.9, 2.4, "Query moves to OS resolver", 13, WHITE, False, 'center')

    ax.annotate('', xy=(5.9, 2.7), xytext=(5.9, 2.7), arrowprops=dict(arrowstyle='->', color=LGRAY, lw=1.5))
    txt(ax, 0.5, 0.4, "Inspect: chrome://net-internals/#dns  |  Firefox: about:networking#dns", 11, ORANGE)
    return fig

# ═══ SLIDE 6 — OS / Hosts File ═══
def s6_os():
    fig, ax = new_fig()
    title_bar(ax, "Step 2: Operating System & Hosts File")
    # OS resolver box
    card(ax, 0.4, 3.8, 5.5, 2.0, ACCENT, 2.5)
    txt(ax, 3.15, 5.3, "OS Resolver", 18, ACCENT, True, 'center')
    txt(ax, 3.15, 4.7, "System-level DNS cache", 14, WHITE, False, 'center')
    txt(ax, 3.15, 4.2, "Shared across all applications", 13, LGRAY, False, 'center')

    # Hosts file box
    card(ax, 6.2, 3.8, 5.4, 2.0, YELLOW, 2.5)
    txt(ax, 8.9, 5.3, "Hosts File", 18, YELLOW, True, 'center')
    txt(ax, 8.9, 4.7, "Static IP-to-hostname mappings", 14, WHITE, False, 'center')
    txt(ax, 8.9, 4.2, "Checked before resolver", 13, LGRAY, False, 'center')

    # Locations
    card(ax, 0.4, 1.0, 11.2, 2.0, ORANGE, 1.5)
    txt(ax, 0.7, 2.5, "Hosts File Locations:", 14, ORANGE, True)
    txt(ax, 0.7, 2.0, "Windows:  C:\\Windows\\System32\\drivers\\etc\\hosts", 13, WHITE)
    txt(ax, 0.7, 1.5, "Linux:    /etc/hosts", 13, WHITE)
    txt(ax, 0.7, 1.0, "macOS:    /etc/hosts", 13, WHITE)

    # Arrow between boxes
    ax.annotate('', xy=(6.2, 4.8), xytext=(5.9, 4.8),
                arrowprops=dict(arrowstyle='->', color=LGRAY, lw=2))

    txt(ax, 0.5, 0.2, "Commands:  ipconfig /displaydns (view)  |  ipconfig /flushdns (clear cache)", 11, ORANGE)
    return fig

# ═══ SLIDE 7 — Recursive Resolver ═══
def s7_resolver():
    fig, ax = new_fig()
    title_bar(ax, "Step 3: Recursive Resolver")
    # Main box
    card(ax, 3.5, 2.8, 5.0, 2.4, ACCENT, 3)
    txt(ax, 6.0, 4.7, "Recursive Resolver", 20, ACCENT, True, 'center')
    txt(ax, 6.0, 4.1, "Does the heavy lifting on your behalf", 14, WHITE, False, 'center')
    txt(ax, 6.0, 3.5, "Caches responses respecting TTL", 13, LGRAY, False, 'center')
    txt(ax, 6.0, 3.0, "Queries root → TLD → authoritative servers", 13, LGRAY, False, 'center')

    # Providers
    providers = [
        ("Google Public DNS", "8.8.8.8  |  8.8.4.4"),
        ("Cloudflare",        "1.1.1.1  |  1.0.0.1"),
        ("Quad9",             "9.9.9.9"),
        ("OpenDNS",           "208.67.222.222  |  208.67.220.220"),
    ]
    for i, (name, ips) in enumerate(providers):
        x = 0.3 + (i % 2) * 6.0
        y = 0.4 + (i // 2) * 0.65
        card(ax, x, y, 5.6, 0.5, GREEN, 1.5)
        txt(ax, x+0.2, y+0.25, name, 12, GREEN, True)
        txt(ax, x+2.5, y+0.25, ips, 11, WHITE)

    # Client arrow
    card(ax, 0.4, 5.3, 2.0, 0.6, YELLOW, 2)
    txt(ax, 1.4, 5.6, "Your Device", 12, YELLOW, True, 'center', 'center')
    ax.annotate('', xy=(3.5, 4.5), xytext=(2.4, 5.6),
                arrowprops=dict(arrowstyle='->', color=LGRAY, lw=2))
    return fig

# ═══ SLIDE 8 — Root Server ═══
def s8_root():
    fig, ax = new_fig()
    title_bar(ax, "Step 4: Root Nameserver")
    # Root box
    card(ax, 3.2, 3.0, 5.6, 2.5, ORANGE, 3)
    txt(ax, 6.0, 5.0, "Root Nameserver", 20, ORANGE, True, 'center')
    txt(ax, 6.0, 4.4, "13 logical authorities (labeled A through M)", 14, WHITE, False, 'center')
    txt(ax, 6.0, 3.8, "Operated by: Verisign, USC/ISI, ICANN, NASA, others", 13, LGRAY, False, 'center')
    txt(ax, 6.0, 3.3, "Anycast-routed — hundreds of physical servers worldwide", 13, LGRAY, False, 'center')

    # Server circles
    letters = 'ABCDEFGHIJKLM'
    for i, l in enumerate(letters):
        angle = 2 * np.pi * i / 13 + 0.2
        r = 0.85
        sx = 6.0 + r * np.cos(angle)
        sy = 4.0 + r * np.sin(angle)
        ax.plot(sx, sy, 'o', color=ACCENT, markersize=10, alpha=0.8)
        txt(ax, sx, sy-0.22, l, 7, WHITE, True, 'center', 'center')

    # Function
    card(ax, 0.4, 0.5, 5.5, 1.5, ACCENT, 2)
    txt(ax, 3.15, 1.5, "What the Root Server Does:", 14, ACCENT, True, 'center')
    txt(ax, 3.15, 1.0, "Does NOT know the IP address", 13, WHITE, False, 'center')
    txt(ax, 3.15, 0.7, "Returns the TLD nameserver for the domain", 13, LGRAY, False, 'center')

    card(ax, 6.2, 0.5, 5.4, 1.5, YELLOW, 2)
    txt(ax, 8.9, 1.5, "Example:", 14, YELLOW, True, 'center')
    txt(ax, 8.9, 1.0, "www.example.com → .com TLD server", 13, WHITE, False, 'center')
    txt(ax, 8.9, 0.7, "Reference: iana.org/domains/root/files", 11, ORANGE, False, 'center')
    return fig

# ═══ SLIDE 9 — TLD Server ═══
def s9_tld():
    fig, ax = new_fig()
    title_bar(ax, "Step 5: TLD Nameserver")
    card(ax, 2.5, 3.2, 7.0, 2.5, ORANGE, 3)
    txt(ax, 6.0, 5.2, "Top-Level Domain (TLD) Nameserver", 20, ORANGE, True, 'center')
    txt(ax, 6.0, 4.5, "Manages a specific TLD (.com, .org, .uk, .dev, etc.)", 14, WHITE, False, 'center')
    txt(ax, 6.0, 3.9, "Operated by domain registries (Verisign, Nominet, Google, etc.)", 13, LGRAY, False, 'center')
    txt(ax, 6.0, 3.5, "Does NOT have the IP → returns the authoritative nameserver instead", 13, LGRAY, False, 'center')

    # TLD types
    tlds = [("gTLDs", ".com .org .net .info"), ("ccTLDs", ".uk .de .jp .in .br"), ("New gTLDs", ".dev .app .blog .io")]
    for i, (ttype, examples) in enumerate(tlds):
        x = 0.3 + i * 4.0
        card(ax, x, 0.5, 3.7, 1.5, GREEN if i == 0 else YELLOW if i == 1 else ACCENT, 2)
        txt(ax, x+1.85, 1.5, ttype, 13, GREEN if i == 0 else YELLOW if i == 1 else ACCENT, True, 'center')
        txt(ax, x+1.85, 1.0, examples, 12, WHITE, False, 'center')

    ax.annotate('', xy=(6.0, 3.2), xytext=(6.0, 2.3),
                arrowprops=dict(arrowstyle='->', color=ACCENT, lw=2.5))
    txt(ax, 6.7, 2.6, "→ Authoritative Server", 12, ACCENT, True)
    return fig

# ═══ SLIDE 10 — Authoritative ═══
def s10_auth():
    fig, ax = new_fig()
    title_bar(ax, "Step 6: Authoritative Nameserver")
    card(ax, 0.4, 2.5, 5.5, 3.0, GREEN, 3)
    txt(ax, 3.15, 5.0, "Authoritative Nameserver", 18, GREEN, True, 'center')
    txt(ax, 3.15, 4.4, "The final stop in the DNS lookup", 14, WHITE, False, 'center')
    txt(ax, 3.15, 3.8, "The server of authority — holds", 13, LGRAY, False, 'center')
    txt(ax, 3.15, 3.4, "the actual DNS resource records", 13, LGRAY, False, 'center')
    txt(ax, 3.15, 2.8, "Returns the final IP address", 13, LGRAY, False, 'center')

    # Records
    card(ax, 6.2, 4.0, 5.4, 1.5, ACCENT, 2)
    txt(ax, 8.9, 5.0, "Common Providers", 14, ACCENT, True, 'center')
    txt(ax, 8.9, 4.5, "Cloudflare DNS  |  AWS Route 53  |  Google Cloud DNS", 12, WHITE, False, 'center')
    txt(ax, 8.9, 4.1, "Azure DNS  |  Namecheap  |  GoDaddy", 12, WHITE, False, 'center')

    # Example records
    records = [("A", "192.0.2.1"), ("AAAA", "2001:db8::1"), ("MX", "mail.example.com"), ("CNAME", "www → example.com")]
    for i, (r, v) in enumerate(records):
        x = 6.2 + (i % 2) * 2.7
        y = 2.5 + (i // 2) * 0.6
        card(ax, x, y, 2.5, 0.45, GREEN, 1.2)
        txt(ax, x+0.15, y+0.22, r, 11, GREEN, True)
        txt(ax, x+1.0, y+0.22, v, 10, WHITE)

    # Arrow from resolver
    ax.annotate('', xy=(5.9, 4.0), xytext=(5.9, 4.0),
                arrowprops=dict(arrowstyle='->', color=ACCENT, lw=2))
    txt(ax, 0.5, 1.5, "After receiving the IP:", 13, YELLOW, True)
    txt(ax, 0.5, 1.0, "1. Resolver caches the response (respects TTL)", 12, WHITE)
    txt(ax, 0.5, 0.6, "2. Resolver returns the IP to your browser", 12, WHITE)
    txt(ax, 0.5, 0.2, "3. Browser opens a TCP connection to that IP", 12, WHITE)
    return fig

# ═══ SLIDE 11 — Recursive vs Iterative ═══
def s11_rec_vs_iter():
    fig, ax = new_fig()
    title_bar(ax, "Recursive vs Iterative Query")

    # Recursive
    card(ax, 0.3, 1.5, 5.6, 4.5, GREEN, 3)
    txt(ax, 3.1, 5.5, "Recursive Query", 22, GREEN, True, 'center')
    items_r = [
        "Client asks the resolver to find",
        "the answer completely",
        "",
        "Resolver does ALL the work:",
        "root → TLD → authoritative",
        "",
        "Client receives only",
        "the final answer (or error)",
        "",
        "Used by end-user devices",
    ]
    for i, item in enumerate(items_r):
        txt(ax, 0.6, 4.8 - i * 0.35, item, 12, LGRAY if item == '' else WHITE if ':' in item else LGRAY)

    # Iterative
    card(ax, 6.1, 1.5, 5.6, 4.5, ACCENT, 3)
    txt(ax, 8.9, 5.5, "Iterative Query", 22, ACCENT, True, 'center')
    items_i = [
        "Server responds with the best",
        "answer it can give",
        "",
        "May return a referral to",
        "another server (next step)",
        "",
        "Client follows referrals:",
        '"Ask server X, not me"',
        "",
        "Used between DNS servers",
    ]
    for i, item in enumerate(items_i):
        txt(ax, 6.4, 4.8 - i * 0.35, item, 12, LGRAY if item == '' else WHITE if ':' in item else LGRAY)
    return fig

# ═══ SLIDE 12 — DNS Records ═══
def s12_records():
    fig, ax = new_fig()
    title_bar(ax, "Common DNS Record Types")
    records = [
        ("A", "IPv4 address", "192.0.2.1", ACCENT),
        ("AAAA", "IPv6 address", "2001:db8::1", ACCENT),
        ("CNAME", "Canonical name / alias", "www → example.com", YELLOW),
        ("MX", "Mail exchange server", "mail.example.com", YELLOW),
        ("TXT", "Text / SPF / DKIM verification", "v=spf1 include:_spf.google.com", ORANGE),
        ("NS", "Authoritative nameserver", "ns1.example.com", ORANGE),
        ("SRV", "Service location", "_sip._tcp.example.com", GREEN),
        ("SOA", "Start of Authority", "Admin info, refresh, retry", GREEN),
    ]
    for i, (rtype, desc, example, color) in enumerate(records):
        col = i % 2
        row = i // 2
        x = 0.3 + col * 6.0
        y = H - 1.7 - row * 1.1
        card(ax, x, y-0.2, 5.7, 0.85, color, 2)
        txt(ax, x+0.2, y+0.2, rtype, 16, color, True)
        txt(ax, x+1.3, y+0.2, desc, 12, WHITE)
        txt(ax, x+0.3, y-0.1, f"Example: {example}", 10, LGRAY)
    return fig

# ═══ SLIDE 13 — Caching ═══
def s13_caching():
    fig, ax = new_fig()
    title_bar(ax, "DNS Caching Hierarchy")
    layers = [
        (5.2, "Browser Cache", "In-memory, per-tab, cleared on close", GREEN, 1.5),
        (4.2, "OS Cache", "System-wide cache", ACCENT, 2.0),
        (3.2, "Resolver Cache", "ISP / Public DNS (shared across users)", ORANGE, 2.5),
        (2.2, "Application Cache", "curl, wget, custom applications", YELLOW, 3.0),
        (1.2, "CDN / Proxy Cache", "Cloudflare, Akamai, reverse proxies", RED, 3.5),
    ]
    for y, label, desc, color, w in layers:
        x = W/2 - w/2
        card(ax, x, y, w, 0.75, color, 3)
        txt(ax, W/2, y+0.5, label, 16 if label == "Browser Cache" else 14, color, True, 'center')
        txt(ax, W/2, y+0.2, desc, 11, LGRAY, False, 'center')
        if y < 5.2:
            ax.annotate('', xy=(W/2 + w/2 + 0.2, y+0.6), xytext=(W/2 + w/2 + 0.2, y+1.0),
                        arrowprops=dict(arrowstyle='->', color=LGRAY, lw=1.5))

    txt(ax, W/2, 0.3, "TTL values determine how long records are cached at each level", 12, LGRAY, False, 'center')
    return fig

# ═══ SLIDE 14 — Tools ═══
def s14_tools():
    fig, ax = new_fig()
    title_bar(ax, "DNS Troubleshooting Tools")
    tools = [
        ("nslookup", "Query DNS interactively or inline", "nslookup example.com", ACCENT),
        ("dig", "Detailed DNS query with full output", "dig example.com A +trace", ACCENT),
        ("ping", "Test reachability (resolves DNS first)", "ping google.com", GREEN),
        ("tracert", "Trace network path to destination", "tracert google.com", GREEN),
        ("ipconfig", "Flush the OS DNS cache (Windows)", "ipconfig /flushdns", YELLOW),
        ("systemd-resolve", "Query systemd-resolved cache (Linux)", "systemd-resolve --status", YELLOW),
    ]
    for i, (tool, desc, cmd, color) in enumerate(tools):
        row = i // 2
        col = i % 2
        x = 0.3 + col * 6.0
        y = H - 1.7 - row * 1.5
        card(ax, x, y-0.25, 5.7, 1.2, color, 2)
        txt(ax, x+0.3, y+0.5, tool, 18, color, True)
        txt(ax, x+0.3, y+0.0, desc, 12, WHITE)
        txt(ax, x+0.3, y-0.15, f"$ {cmd}", 11, GREEN)
    return fig

# ═══ SLIDE 15 — Summary ═══
def s15_summary():
    fig, ax = new_fig()
    title_bar(ax, "Key Takeaways")
    items = [
        "DNS translates domain names into IP addresses",
        "Resolution follows a strict 6-step hierarchy",
        "Caching occurs at every level, dramatically reducing latency",
        "TTL values control how long records are cached",
        "Recursive resolvers do the work; iterative queries chain referrals",
        "DNSSEC, DoH, and DoT improve DNS security and privacy",
        "Tools like dig, nslookup, and ping help diagnose issues",
    ]
    for i, item in enumerate(items):
        y = H - 1.5 - i * 0.55
        ax.plot(0.6, y, 's', color=GREEN, markersize=8)
        txt(ax, 1.0, y, item, 15, LGRAY)

    txt(ax, W/2, 0.4, "DNS Resolution — How the Internet Finds Its Way", 12, ACCENT, False, 'center')
    return fig

# ════════════════════════════════════════════
# BUILD & SAVE
# ════════════════════════════════════════════
slides = [
    s1_title,  s2_what,    s3_why,     s4_pipeline,
    s5_browser, s6_os,     s7_resolver, s8_root,
    s9_tld,    s10_auth,   s11_rec_vs_iter, s12_records,
    s13_caching, s14_tools, s15_summary,
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

bg_color = RGBColor(0x1E, 0x1E, 0x2E)

for i, fn in enumerate(slides):
    print(f"Rendering slide {i+1}...")
    fig = fn()
    img_bytes = fig_to_bytes(fig, 130)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Set background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = bg_color
    # Add the image to fill the slide
    slide.shapes.add_picture(img_bytes, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))

output = "C:\\Users\\mozpc\\AppData\\Local\\Temp\\opencode\\DNS_Resolution.pptx"
prs.save(output)
print(f"\nSaved: {output}")
print(f"Size: {os.path.getsize(output) / 1024:.1f} KB")
